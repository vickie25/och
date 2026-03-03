"""
Export service for sponsor dashboard reports.
Supports PDF, CSV, and PowerPoint formats.
"""
import csv
import io
from datetime import datetime
from django.http import HttpResponse
from .models import Sponsor, SponsorCohort
from . import services as sponsor_services

# Optional imports for PDF/PPTX generation
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class SponsorExportService:
    """Service for generating sponsor dashboard exports"""

    @staticmethod
    def generate_csv_export(sponsor: Sponsor, cohort: SponsorCohort) -> HttpResponse:
        """Generate CSV export of dashboard data"""
        # Get dashboard data
        ai_insights = sponsor_services.SponsorAIService.get_dashboard_ai_insights(cohort)

        # Create CSV response
        response = HttpResponse(content_type='text/csv')
        response['Content-Disposition'] = f'attachment; filename="{sponsor.slug}_dashboard_{datetime.now().strftime("%Y%m%d")}.csv"'

        writer = csv.writer(response)

        # Executive Summary
        writer.writerow(['Executive Summary'])
        writer.writerow(['Metric', 'Value'])
        writer.writerow(['Active Students', cohort.students_enrolled])
        writer.writerow(['Completion Rate', f"{cohort.completion_rate}%"])
        writer.writerow([])

        # Track Performance
        writer.writerow(['Track Performance'])
        writer.writerow(['Track', 'Students', 'Completion Rate', 'Avg Time (days)', 'Hires', 'Avg Salary'])
        tracks = ['defender', 'grc', 'innovation', 'leadership', 'offensive']
        for track in tracks:
            # Mock data - would be calculated from real data
            writer.writerow([
                track.title(),
                '25',  # Mock
                '72.1%',  # Mock
                '38',  # Mock
                '12',  # Mock
                'KES 3.2M'  # Mock
            ])
        writer.writerow([])

        # Top Talent
        writer.writerow(['Top Talent'])
        writer.writerow(['Rank', 'Name', 'Email', 'Readiness Score', 'Completion %'])
        for talent in ai_insights['readiness_scores'][:10]:
            writer.writerow([
                talent['cohort_rank'],
                talent['student_name'],
                talent['student_email'],
                talent['readiness_score'],
                f"{talent['completion_percentage']}%"
            ])

        return response

    @staticmethod
    def generate_pdf_export(sponsor: Sponsor, cohort: SponsorCohort) -> HttpResponse:
        """Generate PDF export of dashboard data"""
        if not REPORTLAB_AVAILABLE:
            # Fallback to CSV if PDF libraries not available
            return SponsorExportService.generate_csv_export(sponsor, cohort)

        # Create PDF response
        response = HttpResponse(content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename="{sponsor.slug}_dashboard_{datetime.now().strftime("%Y%m%d")}.pdf"'

        # Create PDF document
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
        )
        story.append(Paragraph(f"{sponsor.name} Dashboard Report", title_style))
        story.append(Paragraph(f"Generated on {datetime.now().strftime('%B %d, %Y')}", styles['Normal']))
        story.append(Spacer(1, 20))

        # Executive Summary
        story.append(Paragraph("Executive Summary", styles['Heading2']))
        ai_insights = sponsor_services.SponsorAIService.get_dashboard_ai_insights(cohort)

        summary_data = [
            ['Metric', 'Value'],
            ['Active Students', str(cohort.students_enrolled)],
            ['Completion Rate', f"{cohort.completion_rate}%"],
            ['Track', cohort.track_slug.title()],
            ['Start Date', cohort.start_date.strftime('%B %Y') if cohort.start_date else 'TBD'],
        ]

        summary_table = Table(summary_data)
        summary_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 14),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(summary_table)
        story.append(Spacer(1, 20))

        # Top Talent
        story.append(Paragraph("Top Talent", styles['Heading2']))
        talent_data = [['Rank', 'Name', 'Readiness Score', 'Completion %']]
        for talent in ai_insights['readiness_scores'][:10]:
            talent_data.append([
                str(talent['cohort_rank']),
                talent['student_name'],
                str(talent['readiness_score']),
                f"{talent['completion_percentage']}%"
            ])

        talent_table = Table(talent_data)
        talent_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.blue),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(talent_table)

        # Build PDF
        doc.build(story)
        pdf = buffer.getvalue()
        buffer.close()

        response.write(pdf)
        return response

    @staticmethod
    def generate_pptx_export(sponsor: Sponsor, cohort: SponsorCohort) -> HttpResponse:
        """Generate PowerPoint export of dashboard data"""
        if not PPTX_AVAILABLE:
            # Fallback to CSV if PPTX libraries not available
            return SponsorExportService.generate_csv_export(sponsor, cohort)

        # Create PowerPoint response
        response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = f'attachment; filename="{sponsor.slug}_dashboard_{datetime.now().strftime("%Y%m%d")}.pptx"'

        # Create presentation
        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = f"{sponsor.name} Dashboard Report"
        subtitle.text = f"Cohort: {cohort.name}\nGenerated: {datetime.now().strftime('%B %d, %Y')}"

        # Executive Summary slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = 'Executive Summary'
        tf = body_shape.text_frame
        tf.text = f'Active Students: {cohort.students_enrolled}'
        p = tf.add_paragraph()
        p.text = f'Completion Rate: {cohort.completion_rate}%'
        p = tf.add_paragraph()
        p.text = f'Track: {cohort.track_slug.title()}'

        # Top Talent slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = 'Top Talent'
        tf = body_shape.text_frame

        ai_insights = sponsor_services.SponsorAIService.get_dashboard_ai_insights(cohort)
        for talent in ai_insights['readiness_scores'][:5]:
            p = tf.add_paragraph()
            p.text = f"#{talent['cohort_rank']}: {talent['student_name']} - {talent['readiness_score']} readiness"

        # Save to response
        buffer = io.BytesIO()
        prs.save(buffer)
        pptx_data = buffer.getvalue()
        buffer.close()

        response.write(pptx_data)
        return response
